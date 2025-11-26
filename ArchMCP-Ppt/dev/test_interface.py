#!/usr/bin/env python3
"""
Simple test interface for PPT generation modes
"""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent / 'src'))

from bedrock_analyzer import BedrockAnalyzer
from enhanced_search import EnhancedSearch
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import json
import subprocess
import time
import requests
import signal
import atexit

# Global server process
server_process = None

def cleanup_server():
    """Stop server on exit"""
    global server_process
    if server_process:
        server_process.terminate()
        server_process.wait()

atexit.register(cleanup_server)

def start_mcp_server():
    """Start MCP server if not running"""
    global server_process
    
    # Check if server already running
    try:
        response = requests.get('http://localhost:8765', timeout=1)
        print("✓ MCP server already running")
        return True
    except:
        pass
    
    # Start server
    print("🚀 Starting MCP server...")
    server_process = subprocess.Popen(
        [sys.executable, 'src/server.py'],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE
    )
    
    # Wait for server to start
    for i in range(10):
        try:
            requests.get('http://localhost:8765', timeout=1)
            print("✓ MCP server started")
            return True
        except:
            time.sleep(0.5)
    
    print("❌ Failed to start MCP server")
    return False

def test_mcp_server(mode, input_text):
    """Test via MCP server"""
    print(f"\n🌐 MCP SERVER MODE ({mode}): {input_text}")
    
    if not start_mcp_server():
        return None
    
    # Call MCP tool via JSON-RPC
    import json
    
    payload = {
        "jsonrpc": "2.0",
        "id": 1,
        "method": "tools/call",
        "params": {
            "name": "generate_ppt",
            "arguments": {
                "input": input_text,
                "mode": mode
            }
        }
    }
    
    try:
        # MCP uses stdio, so we'll simulate by importing directly
        # In real MCP, this would go through the protocol
        from server import generate_ppt_tool
        result = generate_ppt_tool(input_text, mode)
        
        print(f"✅ Result: {result}")
        return result
    except Exception as e:
        print(f"❌ Error: {e}")
        return None


def generate_ppt_keywords(keywords_str):
    """Test keywords mode"""
    print(f"\n🔧 KEYWORDS MODE: {keywords_str}")
    
    search = EnhancedSearch()
    keywords = [k.strip() for k in keywords_str.split(',')]
    
    icons = []
    for keyword in keywords:
        matches = search.find_service_matches(keyword, {})
        if matches:
            icons.append(matches[0])
            print(f"  ✓ {keyword} → {matches[0]['name']}")
        else:
            print(f"  ✗ {keyword} → No match")
    
    if icons:
        output = create_ppt(icons, f"keywords_{len(icons)}")
        print(f"✅ Created: {output}")
        return output
    return None

def generate_ppt_description(description):
    """Test description mode (AI)"""
    print(f"\n🤖 DESCRIPTION MODE: {description}")
    
    analyzer = BedrockAnalyzer()
    result = analyzer.analyze_architecture(description)
    
    if result and 'services' in result:
        print(f"  AI extracted: {result['services']}")
        
        search = EnhancedSearch()
        icons = []
        for service in result['services']:
            matches = search.find_service_matches(service, {})
            if matches:
                icons.append(matches[0])
                print(f"  ✓ {service} → {matches[0]['name']}")
        
        if icons:
            output = create_ppt(icons, f"description_{len(icons)}")
            print(f"✅ Created: {output}")
            return output
    
    print("❌ No services extracted")
    return None

def create_ppt(icons, name):
    """Create PPT with icons"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    cols = 3
    icon_size = Inches(1.5)
    spacing = Inches(0.5)
    start_x = Inches(1)
    start_y = Inches(1)
    
    for idx, icon in enumerate(icons):
        row = idx // cols
        col = idx % cols
        
        x = start_x + col * (icon_size + spacing)
        y = start_y + row * (icon_size + spacing)
        
        if Path(icon['path']).exists():
            slide.shapes.add_picture(icon['path'], x, y, width=icon_size, height=icon_size)
    
    output = Path('outputs') / f'test_{name}.pptx'
    output.parent.mkdir(exist_ok=True)
    prs.save(str(output))
    return output

def main():
    print("=" * 60)
    print("AWS Architecture PPT Generator - Test Interface")
    print("=" * 60)
    
    while True:
        print("\n📋 Choose mode:")
        print("  1. Keywords mode (direct)")
        print("  2. Description mode (direct AI)")
        print("  3. Keywords via MCP server")
        print("  4. Description via MCP server")
        print("  5. Exit")
        
        choice = input("\nChoice (1-5): ").strip()
        
        if choice == '1':
            keywords = input("Enter keywords (comma-separated): ").strip()
            if keywords:
                generate_ppt_keywords(keywords)
        
        elif choice == '2':
            description = input("Enter architecture description: ").strip()
            if description:
                generate_ppt_description(description)
        
        elif choice == '3':
            keywords = input("Enter keywords (comma-separated): ").strip()
            if keywords:
                test_mcp_server('keywords', keywords)
        
        elif choice == '4':
            description = input("Enter architecture description: ").strip()
            if description:
                test_mcp_server('description', description)
        
        elif choice == '5':
            print("👋 Goodbye!")
            cleanup_server()
            break
        
        else:
            print("❌ Invalid choice")

if __name__ == "__main__":
    main()
