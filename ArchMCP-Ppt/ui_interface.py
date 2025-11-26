#!/usr/bin/env python3
"""
Web UI for PPT generation with icon preview
"""
import sys
from pathlib import Path

# Add ArchMCP-Common to path
sys.path.insert(0, str(Path(__file__).parent.parent / 'ArchMCP-Common'))

from flask import Flask, render_template, request, jsonify, send_file
from enhanced_search import EnhancedSearch
from bedrock_analyzer import BedrockAnalyzer
from pptx import Presentation
from pptx.util import Inches, Pt
import base64
import yaml

app = Flask(__name__, template_folder='templates')
search = EnhancedSearch()

# Load config
config_path = Path(__file__).parent.parent / 'ArchMCP-Common' / 'config' / 'bedrock_config.yaml'
with open(config_path, 'r') as f:
    config = yaml.safe_load(f)

analyzer = BedrockAnalyzer(
    region_name=config['bedrock']['region'],
    model_id=config['bedrock']['model_id']
)
icon_selection_mode = config['bedrock'].get('icon_selection_mode', 'local')
top_k_icons = config['bedrock'].get('top_k_icons', 3)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/search', methods=['POST'])
def search_icons():
    data = request.json
    mode = data.get('mode', 'keywords')
    input_text = data.get('input', '')
    
    icons = []
    debug_info = {'services': [], 'matches': [], 'ambiguous': []}
    
    if mode == 'keywords':
        # Use Bedrock to intelligently select icons from mapping
        keywords = [k.strip() for k in input_text.split(',')]
        debug_info['services'] = keywords
        
        for keyword in keywords:
            # Get all available icons from mapping
            available_icons = list(search.keywords_mapping.keys())
            
            # Use Bedrock to select top K most relevant icons
            selected_icons = analyzer.select_icons_for_keyword(keyword, available_icons, top_n=top_k_icons)
            debug_info['bedrock_response'] = getattr(analyzer, 'last_bedrock_response', 'N/A')
            
            if selected_icons:
                debug_info['matches'].append({
                    'keyword': keyword,
                    'found': selected_icons,
                    'total_matches': len(selected_icons)
                })
                
                # Add icons with metadata
                for icon_name in selected_icons:
                    if icon_name in search.keywords_mapping:
                        icons.append({
                            'name': icon_name,
                            'page': search.keywords_mapping[icon_name].get('page', 0),
                            'category': search.keywords_mapping[icon_name].get('category', ''),
                            'path': search.keywords_mapping[icon_name].get('path', ''),
                            'score': 100
                        })
    else:
        # Use Bedrock to extract services
        if icon_selection_mode == 'llm':
            # LLM selects icons directly
            icon_names = analyzer.analyze_text_description(
                input_text, 
                mode='llm', 
                top_k=top_k_icons,
                available_icons=search.keywords_mapping
            )
            print(f"🔍 Analyzer returned: {icon_names}")
            print(f"🔍 Type: {type(icon_names)}")
            debug_info['bedrock_response'] = getattr(analyzer, 'last_bedrock_response', 'N/A')
            debug_info['mode'] = 'LLM icon selection'
            debug_info['services'] = icon_names if icon_names else []
            
            # Get icon data from mapping
            for icon_name in icon_names:
                # Try with and without .png extension
                lookup_name = icon_name
                if lookup_name not in search.keywords_mapping:
                    lookup_name = icon_name.replace('.png', '') if icon_name.endswith('.png') else icon_name + '.png'
                
                if lookup_name in search.keywords_mapping:
                    icons.append({
                        'name': lookup_name,
                        'page': search.keywords_mapping[lookup_name].get('page', 0),
                        'category': search.keywords_mapping[lookup_name].get('category', ''),
                        'path': search.keywords_mapping[lookup_name].get('path', ''),
                        'score': 100
                    })
                    debug_info['matches'].append({
                        'keyword': icon_name,
                        'found': lookup_name,
                        'score': 100,
                        'total_matches': 1
                    })
        else:
            # Local semantic search
            services = analyzer.analyze_text_description(input_text, mode='local')
            debug_info['bedrock_response'] = getattr(analyzer, 'last_bedrock_response', 'N/A')
            debug_info['replacements'] = getattr(analyzer, 'last_replacements', [])
            debug_info['mode'] = 'Local semantic search'
            debug_info['services'] = services
            for service in services:
                matches = search.find_service_matches(service, {})
                
                # Check if too many matches (ambiguous)
                if len(matches) > 5:
                    debug_info['ambiguous'].append({
                        'keyword': service,
                        'count': len(matches),
                        'options': [m['name'] for m in matches[:10]]
                    })
                
                if matches:
                    # Take top K matches per service
                    top_matches = matches[:top_k_icons]
                    debug_info['matches'].append({
                        'keyword': service,
                        'found': [m['name'] for m in top_matches],
                        'score': matches[0]['score'],
                        'total_matches': len(matches)
                    })
                    icons.extend(top_matches)
    
    # Remove duplicates
    unique_icons = []
    seen_names = set()
    for icon in icons:
        if icon['name'] not in seen_names:
            unique_icons.append(icon)
            seen_names.add(icon['name'])
    
    # Convert icons to base64 for preview
    icons_data = []
    for icon in unique_icons:
        icon_name = icon["name"]
        
        # Use path from icon data if available, otherwise construct from page
        if icon.get('path'):
            possible_paths = [Path(__file__).parent / icon['path']]
        else:
            page = icon.get('page', 0)
            base_path = Path(__file__).parent
            possible_paths = [
                base_path / f'icons/page{page}_icons/{icon_name}',
                base_path / f'icons/page{page}_icons/{icon_name}.png',
                base_path / f'icons/page{page}_icons/{icon_name}.png.png',
            ]
        
        icon_path = None
        for path in possible_paths:
            if path.exists():
                icon_path = path
                break
        
        # Extract page from path if page is 0
        page = icon.get('page', 0)
        if page == 0 and icon.get('path'):
            import re
            match = re.search(r'page(\d+)', icon['path'])
            if match:
                page = int(match.group(1))
        
        if icon_path:
            with open(icon_path, 'rb') as f:
                img_data = base64.b64encode(f.read()).decode()
                icons_data.append({
                    'name': icon_name,
                    'category': icon.get('category', ''),
                    'page': page,
                    'image': f"data:image/png;base64,{img_data}"
                })
        else:
            # Icon file not found, add without image
            icons_data.append({
                'name': icon_name,
                'category': icon.get('category', ''),
                'page': page,
                'image': None
            })
    
    return jsonify({'icons': icons_data, 'debug': debug_info})

@app.route('/generate', methods=['POST'])
def generate_ppt():
    data = request.json
    icon_names = data.get('icons', [])
    
    # Get full icon data
    icons = []
    for name in icon_names:
        if name in search.keywords_mapping:
            mapping = search.keywords_mapping[name]
            
            # Use path from mapping if available, otherwise construct from page
            if mapping.get('path'):
                icon_path = Path(__file__).parent / mapping['path']
            else:
                page = mapping.get('page', 0)
                icon_path = Path(__file__).parent / f'icons/page{page}_icons/{name}'
            
            if icon_path.exists():
                icons.append({
                    'name': name,
                    'path': str(icon_path),
                    'category': mapping.get('category', '')
                })
    
    # Create PPT
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    cols = 3
    icon_size = Inches(1.5)
    spacing = Inches(0.5)
    start_x = Inches(1)
    start_y = Inches(1)
    text_height = Inches(0.4)
    
    for idx, icon in enumerate(icons):
        row = idx // cols
        col = idx % cols
        x = start_x + col * (icon_size + spacing)
        y = start_y + row * (icon_size + spacing + text_height)
        
        if Path(icon['path']).exists():
            # Add icon image
            slide.shapes.add_picture(icon['path'], x, y, width=icon_size, height=icon_size)
            
            # Add text box below icon with name
            text_box = slide.shapes.add_textbox(x, y + icon_size, icon_size, text_height)
            text_frame = text_box.text_frame
            text_frame.text = icon['name'].replace('.png', '').replace('_', ' ')
            text_frame.word_wrap = True
            
            # Format text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(8)
            paragraph.alignment = 1  # Center alignment
    
    output = Path(__file__).parent / 'outputs' / f'ui_generated_{len(icons)}.pptx'
    output.parent.mkdir(exist_ok=True)
    prs.save(str(output))
    
    return jsonify({'file': str(output), 'url': f'/download/{output.name}'})

@app.route('/download/<filename>')
def download(filename):
    output_path = Path(__file__).parent / 'outputs' / filename
    return send_file(output_path, as_attachment=True)

if __name__ == '__main__':
    # Create templates directory
    Path('templates').mkdir(exist_ok=True)
    
    # Create HTML template
    html = '''<!DOCTYPE html>
<html>
<head>
    <title>AWS Architecture PPT Generator</title>
    <style>
        body { font-family: Arial; max-width: 1200px; margin: 50px auto; padding: 20px; }
        h1 { color: #232f3e; }
        .mode-selector { margin: 20px 0; }
        .mode-selector button { padding: 10px 20px; margin-right: 10px; cursor: pointer; }
        .mode-selector button.active { background: #ff9900; color: white; border: none; }
        input, textarea { width: 100%; padding: 10px; margin: 10px 0; font-size: 16px; }
        button { padding: 10px 20px; background: #232f3e; color: white; border: none; cursor: pointer; }
        button:hover { background: #ff9900; }
        .icons-preview { display: grid; grid-template-columns: repeat(auto-fill, minmax(150px, 1fr)); gap: 20px; margin: 20px 0; }
        .icon-card { border: 1px solid #ddd; padding: 10px; text-align: center; }
        .icon-card img { width: 100px; height: 100px; object-fit: contain; }
        .icon-card .name { font-size: 12px; margin-top: 10px; word-break: break-word; }
        .generate-btn { background: #ff9900; font-size: 18px; padding: 15px 30px; margin: 20px 0; }
        .result { margin: 20px 0; padding: 15px; background: #d4edda; border: 1px solid #c3e6cb; }
        .debug { margin: 20px 0; padding: 15px; background: #f8f9fa; border: 1px solid #dee2e6; font-family: monospace; font-size: 12px; }
        .debug h3 { margin-top: 0; font-family: Arial; }
        .ambiguous { margin: 20px 0; padding: 15px; background: #fff3cd; border: 1px solid #ffc107; }
        .ambiguous h3 { margin-top: 0; color: #856404; }
        .ambiguous ul { margin: 10px 0; }
        .spinner {
            border: 4px solid #f3f3f3;
            border-top: 4px solid #ff9900;
            border-radius: 50%;
            width: 50px;
            height: 50px;
            animation: spin 1s linear infinite;
            margin: 0 auto;
        }
        @keyframes spin {
            0% { transform: rotate(0deg); }
            100% { transform: rotate(360deg); }
        }
    </style>
</head>
<body>
    <h1>🏗️ AWS Architecture PPT Generator</h1>
    
    <div class="mode-selector">
        <button id="keywords-btn" class="active" onclick="setMode('keywords')">Keywords Mode</button>
        <button id="description-btn" onclick="setMode('description')">Description Mode (AI)</button>
    </div>
    
    <div id="keywords-input">
        <input type="text" id="keywords" placeholder="Enter keywords: EC2, S3, Lambda" />
    </div>
    
    <div id="description-input" style="display:none;">
        <textarea id="description" rows="3" placeholder="Describe your architecture: Create a web app with EC2 instances behind a load balancer..."></textarea>
    </div>
    
    <button onclick="searchIcons()">🔍 Search Icons</button>
    
    <div id="debug-info" class="debug" style="display:none;"></div>
    
    <div id="icons-preview" class="icons-preview"></div>
    
    <button id="generate-btn" class="generate-btn" style="display:none;" onclick="generatePPT()">📊 Generate PowerPoint</button>
    
    <div id="result" class="result" style="display:none;"></div>
    
    <script>
        let currentMode = 'keywords';
        let foundIcons = [];
        
        function setMode(mode) {
            currentMode = mode;
            document.getElementById('keywords-btn').classList.toggle('active', mode === 'keywords');
            document.getElementById('description-btn').classList.toggle('active', mode === 'description');
            document.getElementById('keywords-input').style.display = mode === 'keywords' ? 'block' : 'none';
            document.getElementById('description-input').style.display = mode === 'description' ? 'block' : 'none';
        }
        
        async function searchIcons() {
            const input = currentMode === 'keywords' 
                ? document.getElementById('keywords').value 
                : document.getElementById('description').value;
            
            if (!input) return;
            
            // Show loading spinner
            const preview = document.getElementById('icons-preview');
            preview.innerHTML = `
                <div style="text-align:center;padding:50px;">
                    <div class="spinner"></div>
                    <p style="margin-top:20px;color:#666;">Analyzing with Bedrock AI...</p>
                </div>
            `;
            document.getElementById('debug-info').innerHTML = '';
            
            try {
                const response = await fetch('/search', {
                    method: 'POST',
                    headers: {'Content-Type': 'application/json'},
                    body: JSON.stringify({mode: currentMode, input: input})
                });
                
                const data = await response.json();
                foundIcons = data.icons;
                
                // Clear loading spinner
                preview.innerHTML = '';
            
            // Show ambiguous warnings
            if (data.debug.ambiguous && data.debug.ambiguous.length > 0) {
                let ambiguousHtml = '<div class="ambiguous">';
                ambiguousHtml += '<h3>⚠️ Ambiguous Terms - Please Be More Specific</h3>';
                data.debug.ambiguous.forEach(amb => {
                    ambiguousHtml += `<p><strong>"${amb.keyword}"</strong> matched ${amb.count} icons. Did you mean:</p>`;
                    ambiguousHtml += '<ul>';
                    amb.options.forEach(opt => {
                        // Clean up the name: remove .png and replace underscores with spaces
                        const cleanName = opt.replace('.png', '').replace(/_/g, ' ');
                        ambiguousHtml += `<li>${cleanName}</li>`;
                    });
                    ambiguousHtml += '</ul>';
                });
                ambiguousHtml += '<p><em>💡 Tip: Use service names (e.g., "Simple Storage Service", "S3 Glacier") instead of technical icon names.</em></p>';
                ambiguousHtml += '<p><em>For S3 specifically, try: "S3 bucket", "S3 Glacier", "S3 Intelligent-Tiering", etc.</em></p>';
                ambiguousHtml += '</div>';
                document.getElementById('icons-preview').insertAdjacentHTML('beforebegin', ambiguousHtml);
            }
            
            // Show debug info
            const debugDiv = document.getElementById('debug-info');
            if (data.debug) {
                let debugHtml = '<h3>🔍 Debug Info</h3>';
                
                if (data.debug.bedrock_response) {
                    debugHtml += '<strong>🤖 Bedrock raw response:</strong><br>';
                    debugHtml += '<pre style="background:#fff;padding:10px;border:1px solid #ddd;">' + data.debug.bedrock_response + '</pre>';
                }
                
                if (data.debug.replacements && data.debug.replacements.length > 0) {
                    debugHtml += '<strong>⚠️ Auto-corrections:</strong><br>';
                    data.debug.replacements.forEach(r => {
                        debugHtml += r + '<br>';
                    });
                    debugHtml += '<br>';
                }
                
                debugHtml += '<strong>Services extracted:</strong> ' + JSON.stringify(data.debug.services) + '<br><br>';
                debugHtml += '<strong>Matches:</strong><br>';
                data.debug.matches.forEach(m => {
                    debugHtml += `"${m.keyword}" → ${m.found} (score: ${m.score}, total: ${m.total_matches})<br>`;
                });
                debugDiv.innerHTML = debugHtml;
                debugDiv.style.display = 'block';
            }
            
            // Reuse preview variable, don't redeclare
            preview.innerHTML = '';
            
            data.icons.forEach(icon => {
                const card = document.createElement('div');
                card.className = 'icon-card';
                const imgHtml = icon.image ? `<img src="${icon.image}" />` : '<div style="padding:50px;background:#f0f0f0;text-align:center;">No Image</div>';
                card.innerHTML = `
                    ${imgHtml}
                    <div class="name">${icon.name}</div>
                    <div style="font-size:10px;color:#666;">Page: ${icon.page || 'N/A'}</div>
                `;
                preview.appendChild(card);
            });
            
            document.getElementById('generate-btn').style.display = data.icons.length > 0 ? 'block' : 'none';
            
            } catch (error) {
                // Hide spinner and show error
                preview.innerHTML = `
                    <div style="text-align:center;padding:50px;color:#d9534f;">
                        <p>❌ Error: ${error.message}</p>
                        <p>Please try again or check the console for details.</p>
                    </div>
                `;
                console.error('Search error:', error);
            }
        }
        
        async function generatePPT() {
            const iconNames = foundIcons.map(i => i.name);
            
            const response = await fetch('/generate', {
                method: 'POST',
                headers: {'Content-Type': 'application/json'},
                body: JSON.stringify({icons: iconNames})
            });
            
            const data = await response.json();
            
            const result = document.getElementById('result');
            result.style.display = 'block';
            result.innerHTML = `✅ PowerPoint created! <a href="${data.url}" download>Download ${data.file.split('/').pop()}</a>`;
        }
    </script>
</body>
</html>'''
    
    with open('templates/index.html', 'w') as f:
        f.write(html)
    
    print("🌐 Starting UI at http://localhost:5000")
    app.run(debug=True, port=5000)
