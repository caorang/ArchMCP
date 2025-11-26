#!/usr/bin/env python3
"""MCP Server for AWS Architecture PowerPoint Generation"""
import json
import asyncio
import sys
from pathlib import Path
from datetime import datetime

# Add ArchMCP-Common to path
sys.path.insert(0, str(Path(__file__).parent.parent / 'ArchMCP-Common'))

from mcp.server import Server
from mcp.types import Tool, TextContent
from mcp_debug_logger import get_debug_logger

app = Server("aws-arch-ppt")

# Paths
ICONS_PATH = Path(__file__).parent / "icons"
KEYWORDS_PATH = Path(__file__).parent / "config" / "keywords_mapping.json"
OUTPUTS_PATH = Path(__file__).parent / "outputs"

# Lazy-loaded modules
_bedrock_analyzer = None
_enhanced_search = None
_pptx_modules = None
_file_server_url = None
_debug_logger = None
_config = None

def get_bedrock_analyzer():
    global _bedrock_analyzer
    if _bedrock_analyzer is None:
        from bedrock_analyzer import BedrockAnalyzer
        import yaml
        
        config_path = Path(__file__).parent.parent / 'ArchMCP-Common' / 'config' / 'bedrock_config.yaml'
        with open(config_path, 'r') as f:
            config = yaml.safe_load(f)
        
        _bedrock_analyzer = BedrockAnalyzer(
            region_name=config['bedrock']['region'],
            model_id=config['bedrock']['model_id']
        )
    return _bedrock_analyzer

def get_config():
    global _config
    if _config is None:
        import yaml
        config_path = Path(__file__).parent.parent / 'ArchMCP-Common' / 'config' / 'bedrock_config.yaml'
        with open(config_path, 'r') as f:
            _config = yaml.safe_load(f)
    return _config

def get_debug_logger_instance():
    global _debug_logger
    if _debug_logger is None:
        config = get_config()
        _debug_logger = get_debug_logger(config)
    return _debug_logger

def get_enhanced_search():
    global _enhanced_search
    if _enhanced_search is None:
        from enhanced_search import EnhancedSearch
        _enhanced_search = EnhancedSearch()
    return _enhanced_search

def get_pptx_modules():
    global _pptx_modules
    if _pptx_modules is None:
        from pptx import Presentation
        from pptx.util import Inches
        _pptx_modules = (Presentation, Inches)
    return _pptx_modules

def get_file_server_url():
    global _file_server_url
    if _file_server_url is None:
        from file_server import start_background
        _file_server_url = start_background()
    return _file_server_url

@app.list_tools()
async def list_tools() -> list[Tool]:
    # Load keywords mapping to report count
    search = get_enhanced_search()
    icon_count = len(search.keywords_mapping)
    
    return [
        Tool(
            name="generate_ppt",
            description=f"Generate a PowerPoint presentation with AWS service icons for the specified architecture ({icon_count} total)",
            inputSchema={
                "type": "object",
                "properties": {
                    "input": {"type": "string", "description": "Architecture description or comma-separated keywords"},
                    "mode": {"type": "string", "enum": ["description", "keywords"], "default": "description"}
                },
                "required": ["input"]
            }
        )
    ]

@app.call_tool()
async def call_tool(name: str, arguments: dict) -> list[TextContent]:
    if name == "generate_ppt":
        logger = get_debug_logger_instance()
        logger.log_mcp_request(name, arguments)
        
        input_text = arguments["input"]
        mode = arguments.get("mode", "description")
        
        try:
            if mode == "description":
                icons = parse_description(input_text)
            else:
                icons = parse_keywords(input_text)
            
            if not icons:
                logger.log_mcp_response([], error="No matching icons found")
                return [TextContent(type="text", text="❌ No matching icons found")]
            
            ppt_path, download_url = create_presentation(icons)
            icon_names = [icon['name'].replace('.png', '') for icon in icons]
            
            logger.log_mcp_response(icons)
            
            return [TextContent(
                type="text",
                text=f"✅ PowerPoint created!\n\n📥 Download: {download_url}\n📁 File: {ppt_path}\n🎨 Icons: {', '.join(icon_names)}"
            )]
        except RuntimeError as e:
            # Bedrock access errors
            error_msg = str(e)
            logger.log_mcp_response([], error=error_msg)
            return [TextContent(type="text", text=error_msg)]
        except Exception as e:
            import traceback
            error_msg = f"❌ Error: {str(e)}\n\n{traceback.format_exc()}"
            logger.log_mcp_response([], error=error_msg)
            return [TextContent(type="text", text=error_msg)]
    
    raise ValueError(f"Unknown tool: {name}")

def parse_description(description: str) -> list[dict]:
    import yaml
    analyzer = get_bedrock_analyzer()
    search = get_enhanced_search()
    
    # Load config
    config_path = Path(__file__).parent.parent / 'ArchMCP-Common' / 'config' / 'bedrock_config.yaml'
    with open(config_path, 'r') as f:
        config = yaml.safe_load(f)
    
    icon_selection_mode = config['bedrock'].get('icon_selection_mode', 'local')
    top_k_icons = config['bedrock'].get('top_k_icons', 3)
    
    if icon_selection_mode == 'llm':
        icon_names = analyzer.analyze_text_description(
            description, 
            mode='llm', 
            top_k=top_k_icons,
            available_icons=search.keywords_mapping
        )
        icons = []
        for icon_name in icon_names:
            lookup_name = icon_name if icon_name in search.keywords_mapping else icon_name.replace('.png', '') if icon_name.endswith('.png') else icon_name + '.png'
            if lookup_name in search.keywords_mapping:
                mapping = search.keywords_mapping[lookup_name]
                icon_path = Path(__file__).parent / (mapping.get('path') or f'icons/page{mapping.get("page", 0)}_icons/{lookup_name}')
                if icon_path.exists():
                    icons.append({'name': lookup_name, 'path': str(icon_path)})
    else:
        services = analyzer.analyze_text_description(description, mode='local')
        icons = []
        for service in services:
            matches = search.find_service_matches(service, {})
            if matches:
                icon_data = matches[0]
                icon_name = icon_data['name']
                if icon_name in search.keywords_mapping:
                    mapping = search.keywords_mapping[icon_name]
                    icon_path = Path(__file__).parent / (mapping.get('path') or f'icons/page{mapping.get("page", 0)}_icons/{icon_name}')
                    if icon_path.exists():
                        icons.append({'name': icon_name, 'path': str(icon_path)})
    return icons

def parse_keywords(keywords: str) -> list[dict]:
    import yaml
    search = get_enhanced_search()
    analyzer = get_bedrock_analyzer()
    
    print(f"🔍 DEBUG: Loaded {len(search.keywords_mapping)} icons in keywords_mapping")
    
    # Load config
    config_path = Path(__file__).parent.parent / 'ArchMCP-Common' / 'config' / 'bedrock_config.yaml'
    with open(config_path, 'r') as f:
        config = yaml.safe_load(f)
    
    top_k_icons = config['bedrock'].get('top_k_icons', 3)
    keywords_list = [k.strip() for k in keywords.split(",")]
    print(f"🔍 DEBUG: Processing keywords: {keywords_list}")
    icons = []
    
    for keyword in keywords_list:
        available_icons = list(search.keywords_mapping.keys())
        print(f"🔍 DEBUG: Available icons count: {len(available_icons)}")
        selected_icons = analyzer.select_icons_for_keyword(keyword, available_icons, top_n=top_k_icons)
        print(f"🔍 DEBUG: Selected icons for '{keyword}': {selected_icons}")
        
        for icon_name in selected_icons:
            if icon_name in search.keywords_mapping:
                mapping = search.keywords_mapping[icon_name]
                icon_path = Path(__file__).parent / (mapping.get('path') or f'icons/page{mapping.get("page", 0)}_icons/{icon_name}')
                print(f"🔍 DEBUG: Checking path: {icon_path}, exists: {icon_path.exists()}")
                if icon_path.exists():
                    icons.append({'name': icon_name, 'path': str(icon_path)})
    
    print(f"🔍 DEBUG: Final icons count: {len(icons)}")
    return icons

def create_presentation(icons: list[dict]) -> tuple[str, str]:
    from pptx.util import Pt
    Presentation, Inches = get_pptx_modules()
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    cols = 3
    icon_size = Inches(1.5)
    spacing = Inches(0.5)
    start_x = Inches(1)
    start_y = Inches(1)
    text_height = Inches(0.4)
    
    for idx, icon in enumerate(icons):
        row = idx // cols
        col = idx % cols
        x = start_x + col * (icon_size + spacing)
        y = start_y + row * (icon_size + spacing + text_height)
        
        icon_path = Path(icon['path'])
        if icon_path.exists():
            # Add icon image
            slide.shapes.add_picture(str(icon_path), x, y, width=icon_size, height=icon_size)
            
            # Add text box below icon
            text_box = slide.shapes.add_textbox(x, y + icon_size, icon_size, text_height)
            text_frame = text_box.text_frame
            text_frame.text = icon['name'].replace('.png', '').replace('_', ' ')
            text_frame.word_wrap = True
            
            # Format text
            paragraph = text_frame.paragraphs[0]
            paragraph.font.size = Pt(8)
            paragraph.alignment = 1  # Center
    
    OUTPUTS_PATH.mkdir(exist_ok=True)
    filename = f"architecture_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    output_path = OUTPUTS_PATH / filename
    prs.save(str(output_path))
    
    download_url = f"{get_file_server_url()}/{filename}"
    return str(output_path), download_url

async def main():
    from mcp.server.stdio import stdio_server
    async with stdio_server() as (read_stream, write_stream):
        await app.run(read_stream, write_stream, app.create_initialization_options())

if __name__ == "__main__":
    asyncio.run(main())
